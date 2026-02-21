import json
import os
import re # For parsing LLM responses and formatting
import math # For number formatting
from pptx import Presentation
from pptx.util import Pt # To set font size
from pptx.dml.color import RGBColor # To set font color
import openpyxl # Assuming fetch_excel_data uses this
import pandas as pd # For saving output CSV
from dotenv import load_dotenv # Assuming F4 uses this
from openai import AzureOpenAI, RateLimitError, APIConnectionError, AuthenticationError, APIError # Assuming F4 uses this
import pprint # Import pprint for better dictionary printing

# ==============================================================================
# ASSUMED EXISTING FUNCTIONS (Ensure these are defined in your script)
# ==============================================================================

def function_1_extract_pptx(slide_number_1_based, prs, current_slide_data):
    """
    Extracts title, relevant text, tables, and identifies charts from a PPTX slide.

    Args:
        slide_number_1_based: The slide number (starting from 1).
        prs: The python-pptx Presentation object.

    Returns:
        A dictionary containing extracted elements, or None if slide fails.
        {
            "slide_number": slide_number_1_based,
            "title": "...",
            "text_content": ["...", "..."], # Combined relevant text
            "tables_data": [[["Row1"], ["Row2"]], [["T2R1"], ["T2R2"]]], # List of tables
            "chart_identifiers": [{"name": "...", "id": ..., "title": "..."}, ...] # Info to find in mapping
        }
    """
    if slide_number_1_based < 1 or slide_number_1_based > len(prs.slides):
        print(f"Error: Slide number {slide_number_1_based} is out of range (1-{len(prs.slides)}).")
        return None

    print(f"\n--- Processing Slide {slide_number_1_based} ---")
    slide = prs.slides[slide_number_1_based - 1]
    # Making a copy of input slide_data_dict
    slide_data = current_slide_data
    slide_data["slide_number_processed"] = slide_number_1_based
    slide_data["title"] = None
    slide_data["text_content"] = [] # Initialize as empty list
    slide_data["tables_data"] = [] # Initialize as empty list
    slide_data["chart_identifiers"] = [] # Initialize as empty list

    # Extract Title (usually the first placeholder)
    try:
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()
            print(f"  Extracted Title: '{slide_data['title']}'")
    except Exception as e:
        print(f"  Warning: Could not extract standard title. {e}")

    # Extract Text, Tables, and Identify Charts
    for shape in slide.shapes:
        # --- Text Extraction (with basic filtering) ---
        if shape.has_text_frame and shape.text.strip():
            # Avoid duplicating title if already extracted
            is_title = slide_data["title"] and shape.has_text_frame and shape.text.strip() == slide_data["title"]
            # Basic filter: Ignore if likely a small element or part of slide master/footer (heuristic)
            # You might need more sophisticated filtering based on position, size, or style
            is_likely_boilerplate = shape.top > (prs.slide_height * 0.85) or shape.height < (prs.slide_height * 0.15)
            if not is_title and not is_likely_boilerplate:
                slide_data["text_content"].append(shape.text.strip())
                print(f"  Extracted Text: '{shape.text.strip()[:50]}...'") # For debugging

        # --- Table Extraction ---
        if shape.has_table:
            print(f"  Found Table: Shape Name='{shape.name}', ID={shape.shape_id}")
            table_data = []
            table = shape.table
            try:
                for row in table.rows:
                    row_data = [cell.text_frame.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                slide_data["tables_data"].append(table_data)
            except Exception as e:
                 print(f"    Warning: Could not fully extract table data for shape '{shape.name}'. Error: {e}")

        # --- Chart Identification ---
        if shape.has_chart:
            chart_title = None
            try:
                if shape.chart.has_title and shape.chart.chart_title.has_text_frame:
                    chart_title = shape.chart.chart_title.text_frame.text.strip()
            except Exception:
                 # Ignore errors getting chart title for now
                 pass
            identifier_info = {"name": shape.name, "id": shape.shape_id, "title": chart_title}
            slide_data["chart_identifiers"].append(identifier_info)
            print(f"  Found Chart: Name='{identifier_info['name']}', ID={identifier_info['id']}, Title='{identifier_info['title']}'")

    print(f"  Finished processing Slide {slide_number_1_based}. Found {len(slide_data['text_content'])} text boxes, {len(slide_data['tables_data'])} tables, {len(slide_data['chart_identifiers'])} charts.")
    return slide_data

def fetch_excel_data(file_path, sheet_name, excel_range=None, range_name=None):
    """
    Fetches data from a specified range or named range in an Excel file using pandas.

    Args:
        file_path (str): Path to the Excel file.
        sheet_name (str): Name of the worksheet.
        excel_range (str, optional): Cell range (e.g., "A1:D10").
        range_name (str, optional): Named range defined in Excel.

    Returns:
        list: A list of lists representing the data (includes header if in range),
              or None if an error occurs. Returns empty list if range is empty.
    """
    if not (excel_range or range_name):
        print(f"  [fetch_excel_data] Error: Provide excel_range or range_name.")
        return None
    if not sheet_name:
        print(f"  [fetch_excel_data] Error: Sheet name not provided.")
        return None
    if not os.path.exists(file_path):
        print(f"  [fetch_excel_data] Error: File not found at '{file_path}'.")
        return None

    target_range = range_name if range_name else excel_range
    print(f"  [fetch_excel_data] Attempting: File='{os.path.basename(file_path)}', Sheet='{sheet_name}', Range='{target_range}'")

    try:
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheet = workbook[sheet_name]

        data_rows = []
        target_cells = None

        if range_name:
            if range_name in workbook.defined_names:
                 dest = workbook.defined_names[range_name].attr_text
                 # Limitation: Simple parsing assuming format like 'SheetName!A1:B5' or just 'A1:B5' if on the target sheet
                 if '!' in dest:
                     ws_title, range_str = dest.split('!', 1)
                     if ws_title != sheet_name:
                         print("  [fetch_excel_data] Warning: Named range points to a different sheet. Ignoring.")
                         workbook.close() # Close workbook to avoid leaks
                 else:
                     range_str = dest
                 target_cells = sheet[range_str]
            else:
                 print(f"  [fetch_excel_data] Error: Named range '{range_name}' not found.")
                 workbook.close()
                 return None
        elif excel_range:
            target_cells = sheet[excel_range]

        if target_cells:
            # Handle single cell vs range
            if isinstance(target_cells, openpyxl.cell.cell.Cell):
                 data_rows.append([target_cells.value])
            else: # It's a tuple of tuples (rows of cells)
                for row_tuple in target_cells:
                    data_rows.append([cell.value for cell in row_tuple])

        workbook.close() # Close workbook when done

        # Convert potential NaNs/NoneTypes if needed for consistency downstream (e.g., replace None with "")
        # data_rows = [[str(cell) if cell is not None else "" for cell in row] for row in data_rows]

        print(f"    Successfully read {len(data_rows)} rows x {len(data_rows[0]) if data_rows else 0} columns.")
        return data_rows # Return list of lists

    except KeyError:
        print(f"  [fetch_excel_data] Error: Sheet '{sheet_name}' not found in '{file_path}'.")
        return None
    except Exception as e:
        print(f"  [fetch_excel_data] Error reading Excel data: {e}")
        return None

# --- Function 2: Add Chart Data from Excel using Mapping ---
def function_2_add_chart_data(slide_number, mapping_data, slide_data):
    # if slide_number == 7:
    #     print("\n\n\n\n\n\n\n\n\n\nHELLLO\n\n\n\n\n\n\n\n\n")
    """
    Finds charts identified in slide_data within the mapping, fetches their
    data from Excel using a helper function, and adds a 'charts_excel_data'
    key to the slide_data dictionary.

    Args:
        slide_number (int): The current slide number (1-based).
        mapping_data (dict): The loaded mapping JSON data.
        slide_data (dict): Dictionary possibly containing 'chart_identifiers' key.
                           This dictionary is modified IN PLACE.
    """
    print(f"\n--- Step 2: Fetching Excel data for Slide {slide_number} charts ---")
    charts_excel_data_list = [] # This will store results for charts successfully processed
    charts_found_on_slide = slide_data.get("chart_identifiers", [])

    if not charts_found_on_slide:
        print("  No chart identifiers found in slide_data (output from Step 1). Skipping Excel fetch.")
        slide_data["charts_excel_data"] = charts_excel_data_list # Ensure key exists, even if empty
        return slide_data # Modify in place, no explicit return needed, but can return slide_data

    # Get default Excel file path from mapping root
    default_excel_file_path = mapping_data.get("excel_path")
    # print(f"Mapping data: \n{mapping_data}\n")
    # print(default_excel_file_path)
    if not default_excel_file_path:
        print("  Warning: No top-level 'excel_path' file path defined in mapping JSON. Cannot fetch Excel data.")
        slide_data["charts_excel_data"] = charts_excel_data_list
        return slide_data

    # Find the mapping rules specific to the current slide number
    slide_mapping_rules = None
    slide_mapping_rules = mapping_data.get("slides", {}).get(str(slide_number))

    if not slide_mapping_rules:
        print(f"  No mapping rules found in JSON for slide {slide_number}. Skipping Excel fetch.")
        slide_data["charts_excel_data"] = charts_excel_data_list
        return slide_data

    mapped_chart_definitions = slide_mapping_rules.get("charts", [])
    if not mapped_chart_definitions:
        print(f"  No 'charts' defined in mapping for slide {slide_number}. Skipping Excel fetch.")
        slide_data["charts_excel_data"] = charts_excel_data_list
        return slide_data

    # --- Match found PPTX charts to mapping definitions and fetch data ---
    print(f"  Comparing {len(charts_found_on_slide)} found charts against {len(mapped_chart_definitions)} mapping definitions.")
    for found_chart in charts_found_on_slide:
        found_chart_name = found_chart.get("name")
        if not found_chart_name:
            print(f"  Skipping PPTX chart (ID: {found_chart.get('id')}, Title: '{found_chart.get('title')}') - lacks shape name.")
            continue

        matched_definition = None
        for definition in mapped_chart_definitions:
            if definition.get("shape_name") == found_chart_name:
                matched_definition = definition
                break

        if matched_definition:
            print(f"  Found mapping for chart '{found_chart_name}'.")
            excel_source_info = matched_definition.get("excel_source", {})
            sheet = excel_source_info.get("sheet")
            # Support both singular 'excel_range' and plural 'excel_ranges' as some mappings use a list
            excel_range = excel_source_info.get("excel_range")
            excel_ranges = excel_source_info.get("excel_ranges")
            range_name = excel_source_info.get("range_name")
            # print(f"\n\n\n\n\n\nMAPPING DATA INCOMING: \n{excel_source_info}\nSheet: {sheet}\nexcel_range: {excel_range}\nexcel_ranges: {excel_ranges}\nrange_name: {range_name}")

            # --- Determine Final Excel File Path ---
            # Check for a specific path within the chart's mapping first
            specific_excel_file = excel_source_info.get("excel_file_path") # Or .get("file") if you use that key
            final_excel_file = specific_excel_file if specific_excel_file else default_excel_file_path

            if not final_excel_file:
                print(f"    Warning: No Excel file path determined for chart '{found_chart_name}' (specific or default). Skipping fetch.")
                continue

            # Must have sheet and at least one range type (support singular 'excel_range', plural 'excel_ranges', or named 'range_name')
            if not sheet or not (excel_range or excel_ranges or range_name):
                print(f"    Warning: Incomplete Excel source details (sheet/range) in mapping for '{found_chart_name}'. Skipping fetch.")
                continue

            # Call the helper function to get data. Support multiple ranges if provided.
            excel_data = None
            try:
                if excel_ranges and isinstance(excel_ranges, (list, tuple)):
                    print(f"    Fetching multiple ranges from: File='{os.path.basename(final_excel_file)}', Sheet='{sheet}', Ranges={excel_ranges}")
                    datasets = []
                    for r in excel_ranges:
                        d = fetch_excel_data(file_path=final_excel_file, sheet_name=sheet, excel_range=r, range_name=None)
                        if d is None:
                            print(f"      Warning: fetch_excel_data returned None for range '{r}'. Skipping this range.")
                            continue
                        datasets.append(d)

                    if datasets:
                        # Combine datasets row-wise by concatenating columns. Pad shorter datasets.
                        max_rows = max(len(d) for d in datasets)
                        col_counts = [len(d[0]) if d and len(d[0])>0 else 0 for d in datasets]
                        combined = []
                        for row_idx in range(max_rows):
                            combined_row = []
                            for ds_idx, ds in enumerate(datasets):
                                if row_idx < len(ds):
                                    row = ds[row_idx]
                                    combined_row += ["" if v is None else v for v in row]
                                else:
                                    # pad with empty strings for missing rows
                                    combined_row += [""] * (col_counts[ds_idx] if col_counts[ds_idx] > 0 else 1)
                            combined.append(combined_row)
                        excel_data = combined
                    else:
                        excel_data = None
                else:
                    # Single range or named range
                    print(f"    Fetching from: File='{os.path.basename(final_excel_file)}', Sheet='{sheet}', Range='{excel_range or range_name}'")
                    excel_data = fetch_excel_data(
                        file_path=final_excel_file,
                        sheet_name=sheet,
                        excel_range=excel_range,
                        range_name=range_name
                    )
            except Exception as e:
                print(f"    ERROR fetching excel ranges: {e}")
                excel_data = None

            if excel_data is not None:
                print(f"    Successfully fetched data for '{found_chart_name}'.")
                chart_output = {
                    "identifier": found_chart_name, # Use shape name as identifier
                    "mapped_title": matched_definition.get("title", found_chart.get("title", "N/A")),
                    "mapped_type": matched_definition.get("chart_type", "N/A"),
                    "data": excel_data # The actual data fetched
                }
                charts_excel_data_list.append(chart_output)
            else:
                print(f"    Failed to fetch or process Excel data for '{found_chart_name}'.")

        else:
            print(f"  Info: PPTX chart '{found_chart_name}' on slide {slide_number} is not defined in the mapping JSON.")

    slide_data["charts_excel_data"] = charts_excel_data_list
    print(f"--- Finished Step 2 (Modified): Added data for {len(charts_excel_data_list)} mapped charts to 'slide_data'. ---")
    return slide_data # Return modified slide_data (though modified in place)

# --- Helper Function to Format List of Lists as Markdown Table ---
def format_as_markdown_table(data_list_of_lists):
    """
    Formats a list of lists (where the first list is headers) into a Markdown table string.

    Args:
        data_list_of_lists (list): A list where the first item is the header row (list of strings),
                                   and subsequent items are data rows (list of strings/numbers/None).

    Returns:
        str: A Markdown formatted table string, or a note if data is empty/invalid.
    """
    if not data_list_of_lists or len(data_list_of_lists) < 1:
        return " (No data available or header missing)\n"

    header = data_list_of_lists[0]
    data_rows = data_list_of_lists[1:]
    num_columns = len(header)

    if num_columns == 0:
        return " (Table header is empty)\n"

    # Create header row string: | Col1 | Col2 | ... |
    header_str = "| " + " | ".join(str(h).strip() for h in header) + " |"
    # Create separator row string: |---|---|...|
    separator_str = "|-" + "-|".join(['---'] * num_columns) + "-|" # Simple separator

    # Create data row strings: | Val1 | Val2 | ... |
    row_strs = []
    for row in data_rows:
        # Ensure row has same number of elements as header for a clean table
        # Replace None with empty string for display
        formatted_row = [str(cell).strip() if cell is not None else "" for cell in row]
        # Pad or truncate row if necessary (simple padding shown)
        if len(formatted_row) < num_columns:
            formatted_row += [""] * (num_columns - len(formatted_row))
        elif len(formatted_row) > num_columns:
            formatted_row = formatted_row[:num_columns] # Truncate

        row_strs.append("| " + " | ".join(formatted_row) + " |")

    # Add a note if there were data rows, otherwise indicate header only
    if not data_rows:
        return "\n".join([header_str, separator_str, "(No data rows)"]) + "\n"
    else:
        return "\n".join([header_str, separator_str] + row_strs) + "\n"

def format_context_data(slide_num, slide_data):
    """ Formats data for ONE slide into a string for the LLM prompt. """
    context_str = ""
    if slide_data:
        # Simplified formatting - use format_as_markdown_table for tables/charts
        context_str += f"\n**Data from Slide {slide_num} ({slide_data.get('title', 'No Title')})**\n"
        if slide_data.get("text_content"):
            context_str += "* Texts:\n" + "\n".join([f"  - {t}" for t in slide_data["text_content"]]) + "\n"
        if slide_data.get("tables_data"):
            context_str += "* Tables:\n"
            for i, tbl in enumerate(slide_data["tables_data"]):
                 context_str += f"  - Table {i+1}:\n{format_as_markdown_table(tbl)}\n"
        if slide_data.get("charts_excel_data"): # Data from F2
            context_str += "* Chart Data (from Linked Excel Source):\n"
            for i, chrt in enumerate(slide_data["charts_excel_data"]):
                 context_str += f"  - Chart {i+1} ({chrt.get('mapped_title', 'N/A')}):\n{format_as_markdown_table(chrt.get('data',[]))}\n"
        # Add extracted chart descriptions if you implement that instead of F2's Excel fetch
    else:
         context_str += f"\n**Data from Slide {slide_num} (Not Available)**\n"
    return context_str

def function_4_call_llm(prompt_string, client, model_id="gpt-4o-mini", temperature=0.3, max_tokens_response=500):
    """
    Sends the formatted prompt to the specified OpenAI Chat Completions model
    and returns the generated text content (insights).

    Args:
        prompt_string (str): The complete prompt generated.
        client (AzureOpenAI): The initialized AzureOpenAI client object.
        model_id (str): The OpenAI model deployment ID.
        temperature (float): Controls creativity/randomness (0.0-2.0).
        max_tokens_response (int): Maximum number of tokens allowed for the response.

    Returns:
        str: The generated text insights from the LLM, or None if an error occurs.
    """
    print("\n--- Step 4: Calling OpenAI API ---")
    if not client:
        print("  Error: Invalid OpenAI client provided to function_4_call_llm.")
        return None
    if not prompt_string or not isinstance(prompt_string, str):
        print("  Error: Invalid prompt string provided. Cannot make API call.")
        return None
    print(f"  Model: '{model_id}', Temperature: {temperature}, Max Tokens: {max_tokens_response}")
    try:
        messages_for_api = [{"role": "user", "content": prompt_string}]
        response = client.chat.completions.create(
            model=model_id,
            messages=messages_for_api,
            temperature=temperature,
            max_tokens=max_tokens_response,
            n=1,
            stop=None
        )
        generated_insights = response.choices[0].message.content.strip()
        usage = response.usage
        print(f"  API Call Success. Usage: Prompt={usage.prompt_tokens}, Completion={usage.completion_tokens}, Total={usage.total_tokens}")
        print("--- Finished Step 4: Received response from LLM successfully. ---")
        return generated_insights
    except AuthenticationError as e: print(f"  FATAL API Error: OpenAI Authentication Failed (Status Code: {e.status_code}). Check API key. Error: {e}"); return None
    except RateLimitError as e: print(f"  API Error: OpenAI Rate Limit Exceeded (Status Code: {e.status_code}). Error: {e}"); return None
    except APIConnectionError as e: print(f"  Network Error: Could not connect to OpenAI API. Error: {e}"); return None
    except APIError as e:
        if e.status_code == 404: print(f"  API Error: Deployment Not Found (Status Code: 404). Check deployment name '{model_id}'. Error: {e}")
        else: print(f"  API Error: OpenAI returned an error (Status Code: {e.status_code}). Error: {e}")
        return None
    except Exception as e:
        print(f"  An unexpected error occurred during the OpenAI API call: {e}")
        import traceback; traceback.print_exc(); return None

# ==============================================================================
# NEW/UPDATED HELPER FUNCTIONS FOR FORMATTING
# ==============================================================================

def format_large_number(num_str):
    """Formats a number string into K/M format if large enough."""
    try:
        num = float(num_str.replace(',', '').replace('K','000').replace('M','000000'))
        if abs(num) >= 1_000_000: return f"{num / 1_000_000:.1f}M"
        elif abs(num) >= 1_000: return f"{num / 1000:.0f}K"
        else:
            if num == math.floor(num): return str(int(num))
            else: return f"{num:.1f}"
    except (ValueError, TypeError): return num_str

def format_percentage(perc_match):
    """
    Formats a percentage string based on magnitude.
    Returns 'a significant increase' or 'a significant decrease' if outside bounds.
    """
    perc_str = perc_match.group(1) # The full matched percentage string (e.g., "+150% WoW")
    num_part_str = perc_match.group(2) # The numeric part (e.g., "+150")
    try:
        num = float(num_part_str)
        if num > 120:
            # Return the replacement noun phrase
            return "a significant increase"
        elif num < -120:
            # Return the replacement noun phrase
            return "a significant decrease"
        else:
            # Return the original matched string if within bounds
            return perc_str
    except (ValueError, TypeError):
        return perc_str # Return original if conversion fails

# ==============================================================================
# UPDATED build_prompt_for_slide FUNCTION (Numbered Lists, Grammar Fix)
# ==============================================================================

def build_prompt_for_slide(slide_num, slide_data, slide_2_context_snippets=None, chart_focus=None):
    """
    Builds the LLM prompt dynamically. Handles chart-specific focus.
    Requests numbered lists with spacing and uses adjusted templates.

    Args:
        slide_num (int): The slide number being processed.
        slide_data (dict): Extracted data for this slide (or just a single chart's data).
        slide_2_context_snippets (dict, optional): Snippets for slide 2.
        chart_focus (dict, optional): Info about the specific chart to focus on.

    Returns:
        str: The constructed prompt string, or None if slide not handled.
    """
    print(f"  Building prompt for Slide {slide_num}" + (f" (Focus: Chart '{chart_focus.get('title', chart_focus.get('name', 'Unknown'))}')" if chart_focus else ""))

    # --- Base Instructions ---
    base_instructions = """You are an expert analyst generating insights for a presentation slide.
Analyze the provided data context specified below.
Generate insights strictly based ONLY on the provided data and follow the specific output requirements."""

    # --- Default values ---
    template = ""
    sub_topic_focus = ""
    extra_instructions = ""
    context_extraction_instruction = ""
    # --- UPDATED Output Requirements ---
    output_requirements = """
**Output Requirements:**
1. Generate the 'Key Highlights' first, following the structure implied by the template below. Generate 3-4 numbered bullet points (e.g., '1.', '2.') with a line break between each point. Use professional, readable full sentences.
2. Fill the `[...]` placeholders within the Key Highlights using specific information derived *solely* from the 'Provided Data Context' or 'Additional Context'. Use 'N/A' if data is missing. For placeholders like '[Change Description]', provide the specific percentage change (e.g., "+15% WoW", "-8%") if available and within +/-120%, otherwise use a qualitative phrase like "a significant increase", "a slight decrease", "minimal change", etc.
3. **CRITICAL:** Ensure the final Key Highlights output text does NOT contain any square brackets `[` or `]` or the heading "Key Highlights:".
4. After the Key Highlights, on a **NEW LINE**, provide the 'Summary Phrase'. Format it EXACTLY like this: `Summary Phrase: [Concise phrase summarizing the most impactful highlight, max 10 words.]`
{context_extraction_instruction}""" # Include context extraction placeholder

    final_request = """**Generate the Key Highlights first (as a numbered list with line breaks), then the Summary Phrase on a new line below it (and the Context Snippet if applicable):**"""

    # --- Handle Chart-Specific Prompts (Slides 7, 10) ---
    if chart_focus:
        chart_name = chart_focus.get('name', 'this chart')
        chart_title = chart_focus.get('title', chart_name)
        base_instructions = f"""You are an expert analyst generating insights for a specific chart ('{chart_title}') on a presentation slide.
Analyze ONLY the provided data context for this chart below. Starting from the latest week, just analyze the previous 4 weeks. Also display week number as last 2 digits instead of using 'K'; eg: for 202506, display just Week 06. Try to avoid commas, semicolons, fullstops and give a complete sentence
Generate 1-2 concise, readable numbered bullet points (e.g., '1.', '2.') summarizing the key findings from ONLY this chart. Ensure a line break between points."""
        sub_topic_focus = f"Key takeaways from chart '{chart_title}'"
        template = """1. [Insight 1 about this specific chart]
2. [Optional Insight 2 about this specific chart]""" # Numbered template
        # Override output requirements
        output_requirements = """
**Output Requirements:**
1. Generate 1-2 numbered bullet points summarizing the key findings from ONLY the provided chart data. Ensure a line break between points.
2. Use professional, readable full sentences, not just phrases.
3. **CRITICAL:** Ensure the final output text does NOT contain any square brackets `[` or `]`. Do not include any heading like "Key Highlights:". """
        final_request = "**Generate the required numbered bullet points now:**"
        context_data_str = f"\n**Data Context (Chart: {chart_title} on Slide {slide_num})**\n"
        context_data_str += format_as_markdown_table(chart_focus.get('data', [])) # Assumed helper

    # --- Handle Full Slide Prompts (Other Slides) ---
    else:
        # --- UPDATED Templates with Numbering and [Change Description] ---
        if slide_num == 2:
            template = """1. App Installs: Installs showed [Change Description] at [Value], driven by [Primary Driver]. [Specific driver context based on snippets below].
2. WAU: WAU showed [Change Description] at [Value].
3. [Add key insight 3]
4. [Optional key insight 4]"""
            sub_topic_focus = "Weekly Performance including App Installs (Organic/Owned/Paid drivers) and WAU."
            extra_instructions = """Synthesize information. For Installs, specify the primary driver (organic/owned/paid). Then add relevant context:
- If organic driven, mention OS: {organic_os}
- If owned driven, mention top campaign: {owned_campaign}
- If paid driven, mention media partner: {paid_partner}
Fill these details using the provided context snippets. State 'N/A' if a snippet is missing or not applicable."""
            snippets = slide_2_context_snippets or {}
            extra_instructions = extra_instructions.format(
                organic_os=snippets.get('organic_os', 'N/A'),
                owned_campaign=snippets.get('owned_campaign', 'N/A'),
                paid_partner=snippets.get('paid_partner', 'N/A')
            )
        elif slide_num == 3:
            template = """1. [Provide key insight about installs breakout, focusing on OS if relevant for organic].
2. [Add key insight 2]
3. [Add key insight 3]"""
            sub_topic_focus = "Installs Breakout details, particularly the Operating System driving organic performance if applicable."
            context_extraction_instruction = "\nVERY IMPORTANT: On a new line *after* the Summary Phrase, identify the primary Operating System (OS) driving organic installs based *only* on the provided data for this slide. Format it EXACTLY like this: 'Context Snippet for Slide 2 - Organic OS: [Detected OS Name or N/A]'."
        elif slide_num == 4:
            template = """1. Overall: Owned installs showed [Change Description] at [Value].
2. Key Driver Highlight: The primary driver for owned installs growth/decline was [Identify main driver/campaign]. [Optional: Add brief detail].
3. [Add key insight 3]
4. [Optional key insight 4]"""
            sub_topic_focus = "Overall Owned Installs performance vs WoW, and identify the primary 2-3 drivers or campaigns influencing growth/decline."
            context_extraction_instruction = "\nVERY IMPORTANT: On a new line *after* the Summary Phrase, identify the specific campaign name witnessing the highest growth for owned installs based *only* on the provided data for this slide. Format it EXACTLY like this: 'Context Snippet for Slide 2 - Owned Campaign: [Campaign Name or N/A]'."
        elif slide_num == 5:
            template = """1. Overall: Total paid installs showed [Change Description] at [Value].
2. Primary Driver: The primary driver for paid spend/installs was [Identify main driver/partner/campaign], driven by [Specific reason if available].
3. [Add key insight 3]
4. [Optional key insight 4]"""
            sub_topic_focus = "Overall Paid Installs performance vs WoW, and identify the primary driver (e.g., Google) and any reason for change (e.g., spend increase)."
            context_extraction_instruction = "\nVERY IMPORTANT: On a new line *after* the Summary Phrase, identify the primary media partner associated with paid installs based *only* on the provided data for this slide. Format it EXACTLY like this: 'Context Snippet for Slide 2 - Paid Partner: [Partner Name or N/A]'."
        elif slide_num == 6:
             template = """1. Sessions/Cost: [Value] sessions at [Cost/Session Value] $/Session and CPA of [CPA Value].
2. Re-engagement: [Value] users re-engaged ([Value]% of WAU).
3. [Add key insight 3]
4. [Optional key insight 4]"""
             sub_topic_focus = "Sessions, cost per session, CPA, and re-engagement rate."
        elif slide_num == 8: # Slide 8 is now standard
            template = """1. Segment Growth: Reactivated users showed [Change Description], New users [Change Description], Current users [Change Description] YoY(compare the current reporting week with last year's same week+1, eg: 202506 vs 2020407; also display week number as last 2 digits instead of using 'K'; eg: for 202506, display just Week 06. Try to avoid commas, semicolons, fullstops and give a complete sentence).
2. [Add key insight 2 about segments]
3. [Add key insight 3 about segments]"""
            sub_topic_focus = "Segment-wise user growth (Reactivated, New, Current) compared YoY."
        # --- Slides 7 and 10 are handled chart-by-chart ---
        elif slide_num in [1, 9, 11]: # Generic for others
             print(f"  Using generic template for Slide {slide_num}.")
             template = """1. [Insight 1]
2. [Insight 2]
3. [Insight 3]
4. [Optional Insight 4]"""
             sub_topic_focus = "General summary of key data points."
        else:
            print(f"  Warning: No specific template logic defined for Slide {slide_num}. Skipping prompt generation.")
            return None

        # Format context data for the whole slide
        context_data_str = format_context_data(slide_num, slide_data) # Assumed helper function
        output_requirements = output_requirements.format(context_extraction_instruction=context_extraction_instruction)

    # --- Construct the Final Prompt ---
    prompt = f"""{base_instructions}

**Analysis Context & Focus:**
- Target Slide: {slide_num} ({slide_data.get('title', 'N/A') if not chart_focus else 'Chart Focus'})
- Key Focus Areas: {sub_topic_focus}
{extra_instructions}
{output_requirements}
"""
    if not chart_focus:
        prompt += f"""
**Template Structure Guide (Follow Strictly for Key Highlights):**
{template}"""

    prompt += f"""
**Provided Data Context:**
{context_data_str}

{final_request}
"""
    # Add start for LLM for numbered lists if template starts that way
    if template.strip().startswith('1.'):
        prompt += "\n1. "

    print(f"  Prompt built for Slide {slide_num}" + (f" (Chart Focus: {chart_focus.get('title', chart_focus.get('name'))})" if chart_focus else ""))
    return prompt

# ==============================================================================
# UPDATED parse_llm_response FUNCTION (No changes needed for these requests)
# ==============================================================================

def parse_llm_response(slide_num, response_text, is_chart_specific=False):
    """
    Parses the LLM response. If not chart_specific, separates Key Highlights,
    Summary Phrase, and context snippet. Applies formatting rules.

    Args:
        slide_num (int): The slide number the response is for.
        response_text (str): The raw text output from the LLM.
        is_chart_specific (bool): Flag indicating if the prompt was for a single chart.

    Returns:
        dict: Contains 'main_insight' (str), and optionally 'summary_phrase', 'context_snippet'.
    """
    if not response_text:
        if is_chart_specific: return {'main_insight': "[Error: No response from LLM]"}
        else: return {'main_insight': "[Error: No response from LLM]", 'summary_phrase': None, 'context_snippet': None}

    main_insight = response_text
    summary_phrase = None
    context_snippet = None

    # --- Extract context snippet first (only if NOT chart-specific) ---
    if not is_chart_specific and slide_num in [3, 4, 5]:
        patterns = {
            3: r"Context Snippet for Slide 2 - Organic OS:\s*(.*)",
            4: r"Context Snippet for Slide 2 - Owned Campaign:\s*(.*)",
            5: r"Context Snippet for Slide 2 - Paid Partner:\s*(.*)"
        }
        pattern = patterns.get(slide_num)
        if pattern:
            match = re.search(pattern, main_insight, re.IGNORECASE | re.MULTILINE | re.DOTALL)
            if match:
                context_snippet = match.group(1).strip(); main_insight = re.sub(pattern, '', main_insight, flags=re.IGNORECASE | re.MULTILINE | re.DOTALL).strip()
                print(f"    Successfully parsed context snippet for Slide {slide_num}: '{context_snippet}'")
            else: print(f"    Warning: Could not parse context snippet for Slide {slide_num}. Check LLM output format.")

    # --- Extract Summary Phrase (only if NOT chart-specific) ---
    if not is_chart_specific:
        summary_pattern = r"Summary Phrase:\s*(.*)"
        summary_match = re.search(summary_pattern, main_insight, re.IGNORECASE | re.MULTILINE | re.DOTALL)
        if summary_match:
            summary_phrase = summary_match.group(1).strip(); main_insight = re.sub(summary_pattern, '', main_insight, flags=re.IGNORECASE | re.MULTILINE | re.DOTALL).strip()
            print(f"    Successfully parsed summary phrase: '{summary_phrase}'")
        else: print(f"    Warning: Could not parse 'Summary Phrase:' from LLM response for slide {slide_num}.")

    # --- Clean up main insight ---
    main_insight = re.sub(r'^```text\s*', '', main_insight, flags=re.IGNORECASE).strip()
    main_insight = re.sub(r'\s*```$', '', main_insight, flags=re.IGNORECASE).strip()
    main_insight = re.sub(r'^Provide a brief summary of.*\.\s*', '', main_insight, flags=re.IGNORECASE).strip()
    main_insight = re.sub(r'^\s*(\*\*?)?Key Highlights:?(\*\*?)?\s*', '', main_insight, flags=re.IGNORECASE | re.MULTILINE).strip()
    main_insight = re.sub(r'^\s*(\*\*?)?[\w\s]+ Insight:?(\*\*?)?\s*', '', main_insight, flags=re.IGNORECASE | re.MULTILINE).strip()
    main_insight = re.sub(r'^\s*(\*\*?)?[\w\s]+ Performance:?(\*\*?)?\s*', '', main_insight, flags=re.IGNORECASE | re.MULTILINE).strip()
    main_insight = re.sub(r'\[[^\]]*\]', 'N/A', main_insight).strip()

    # --- Apply Formatting Rules to main_insight ---
    percent_pattern = r'(([+-]?\d+(?:\.\d+)?)(\s*%\s*(?:WoW)?\b))' # Original pattern to find percentages
    try:
        # Apply the formatting function which now returns phrases for >120%
        main_insight = re.sub(percent_pattern, format_percentage, main_insight, flags=re.IGNORECASE)
        print("    Applied percentage formatting.")
    except Exception as e: print(f"    Warning: Error during percentage formatting: {e}")

    def replace_numbers(text):
        processed_text = text; num_pattern = r'(?<![KM])\b(\d{1,3}(?:,\d{3})+|\d{4,})\b(?![KM])'
        processed_text = re.sub(num_pattern, lambda m: format_large_number(m.group(1)), processed_text); return processed_text
    try: main_insight = replace_numbers(main_insight); print("    Applied number abbreviation formatting.")
    except Exception as e: print(f"    Warning: Error during number abbreviation formatting: {e}")

    main_insight = re.sub(r'\n{2,}', '\n', main_insight).strip()

    # Return appropriate dictionary structure
    if is_chart_specific:
        return {'main_insight': main_insight}
    else:
        return {'main_insight': main_insight, 'summary_phrase': summary_phrase, 'context_snippet': context_snippet}


# ==============================================================================
# UPDATED FUNCTION TO PASTE ALL TEXT TYPES (Numbered Lists & Spacing)
# ==============================================================================

def paste_all_text_to_ppt(prs, all_text_map):
    """
    Pastes various types of text (Key Highlights, Summary Phrase, Chart Specific)
    into specified placeholder shapes in the PowerPoint presentation with appropriate formatting.
    Handles numbered lists and spacing for Key Highlights.

    Args:
        prs (pptx.Presentation): The Presentation object (loaded).
        all_text_map (dict): Dictionary mapping slide numbers (int) to dictionaries
                             where keys are shape names and values are dicts
                             {'text': insight_text, 'type': 'key_highlight' | 'summary_phrase' | 'chart_specific'}.
    """
    print("\n--- Pasting All Generated Text into PowerPoint ---")
    pasted_count = 0
    failed_pastes = []

    # --- Define Formatting ---
    HEADING_TEXT = "Key Highlights"
    HIGHLIGHTS_FONT_NAME = "Poppins"; HIGHLIGHTS_FONT_FALLBACK = "Calibri"
    HEADING_FONT_SIZE_PT = 8; INSIGHT_FONT_SIZE_PT = 7
    HIGHLIGHTS_FONT_COLOR_RGB = RGBColor.from_string("001E60")

    SUMMARY_FONT_NAME = "Poppins"; SUMMARY_FONT_FALLBACK = "Calibri"
    SUMMARY_FONT_SIZE_PT = 9 # Specific size
    SUMMARY_FONT_COLOR_RGB = RGBColor.from_string("001E60")

    CHART_SPECIFIC_FONT_NAME = "Poppins"; CHART_SPECIFIC_FONT_FALLBACK = "Calibri"
    CHART_SPECIFIC_FONT_SIZE_PT = 7 # Same as insight text
    CHART_SPECIFIC_FONT_COLOR_RGB = RGBColor.from_string("001E60")

    # Sort slide numbers to process in order
    for slide_num in sorted(all_text_map.keys()):
        shape_texts = all_text_map[slide_num]
        if not isinstance(shape_texts, dict): continue

        if not (1 <= slide_num <= len(prs.slides)):
            print(f"  Skipping Slide {slide_num}: Slide number out of range.")
            failed_pastes.append(f"Slide {slide_num} (Out of Range)")
            continue

        slide = prs.slides[slide_num - 1]
        print(f"  Processing Slide {slide_num} for pasting...")

        for shape_name, text_info in shape_texts.items():
            text_to_paste = text_info.get('text')
            paste_type = text_info.get('type')

            if not text_to_paste or text_to_paste.startswith("[Error"):
                print(f"    Skipping paste to '{shape_name}': Text missing or error.")
                failed_pastes.append(f"Slide {slide_num} / Shape '{shape_name}' (Text Error)")
                continue

            shape_found = False
            for shape in slide.shapes:
                if hasattr(shape, 'name') and shape.name.strip() == shape_name.strip():
                    if shape.has_text_frame:
                        print(f"    Found shape '{shape_name}'. Pasting text (Type: {paste_type})...")
                        try:
                            tf = shape.text_frame; tf.clear()
                            # Ensure first paragraph exists for adding runs
                            if not tf.paragraphs: tf.add_paragraph()
                            p = tf.paragraphs[0]
                            for r in p.runs: r._r.getparent().remove(r._r) # Clear runs from first para

                            # Apply formatting based on type
                            if paste_type == 'key_highlight' or paste_type == 'chart_specific':
                                # Add Heading to first paragraph
                                run_heading = p.add_run()
                                run_heading.text = HEADING_TEXT; run_heading.font.name = HIGHLIGHTS_FONT_NAME
                                run_heading.font.size = Pt(HEADING_FONT_SIZE_PT); run_heading.font.color.rgb = HIGHLIGHTS_FONT_COLOR_RGB
                                run_heading.font.bold = True

                                # Add Insight Text lines in subsequent paragraphs
                                insight_lines = text_to_paste.split('\n')
                                item_number = 1 # Start numbering
                                for line in insight_lines:
                                    line = line.strip();
                                    if not line: continue # Skip empty lines

                                    p_insight = tf.add_paragraph() # New paragraph for each line/bullet
                                    # Remove potential leading bullet/number from LLM output
                                    text_for_run = re.sub(r'^\s*[\*\-\d]+\.?\s*', '', line).strip()

                                    # Add run with the number and the text
                                    run_insight = p_insight.add_run()
                                    run_insight.text = f"{item_number}. {text_for_run}" # Add number prefix
                                    item_number += 1

                                    # Apply formatting
                                    font_name = HIGHLIGHTS_FONT_NAME if paste_type == 'key_highlight' else CHART_SPECIFIC_FONT_NAME
                                    font_size = INSIGHT_FONT_SIZE_PT if paste_type == 'key_highlight' else CHART_SPECIFIC_FONT_SIZE_PT
                                    font_color = HIGHLIGHTS_FONT_COLOR_RGB if paste_type == 'key_highlight' else CHART_SPECIFIC_FONT_COLOR_RGB
                                    run_insight.font.name = font_name; run_insight.font.size = Pt(font_size)
                                    run_insight.font.color.rgb = font_color; run_insight.font.bold = False
                                    # Ensure spacing between paragraphs (adjust Pt value if needed)
                                    p_insight.space_before = Pt(0)
                                    p_insight.space_after = Pt(3) # Add small space after each bullet point paragraph

                            elif paste_type == 'summary_phrase':
                                run = p.add_run() # Use the first paragraph
                                run.text = text_to_paste
                                run.font.name = SUMMARY_FONT_NAME
                                run.font.size = Pt(SUMMARY_FONT_SIZE_PT)
                                run.font.color.rgb = SUMMARY_FONT_COLOR_RGB
                                run.font.bold = False

                            else: # Default paste if type is unknown
                                 run = p.add_run(); run.text = text_to_paste

                            shape_found = True; pasted_count += 1
                            print(f"    Successfully pasted text to shape '{shape_name}'.")
                            break # Move to next shape_name in map
                        except Exception as e:
                            print(f"    ERROR pasting text into shape '{shape_name}': {e}")
                            failed_pastes.append(f"Slide {slide_num} / Shape '{shape_name}' (Pasting Error)")
                            shape_found = True; break # Stop trying for this shape
                    else:
                        print(f"    Warning: Found shape '{shape_name}', but it has no text frame.")
                        failed_pastes.append(f"Slide {slide_num} / Shape '{shape_name}' (No Text Frame)")
                        shape_found = True; break # Stop trying for this shape

            if not shape_found:
                print(f"    Warning: Target shape '{shape_name}' not found on slide {slide_num}.")
                failed_pastes.append(f"Slide {slide_num} / Shape '{shape_name}' (Not Found)")

    print(f"--- Finished Pasting All Text: Pasted {pasted_count} text blocks. Failures/Skips: {failed_pastes} ---")


# ==============================================================================
# MAIN LOGIC FUNCTION (UPDATED to store prompts)
# ==============================================================================

def generate_and_paste_insights(config_file_path):
    """
    Main function to generate insights, handle chart-specific cases,
    paste into PowerPoint, store prompts, and save.

    Args:
        config_file_path (str): Path to the unified configuration JSON file.
    """
    print(f"Starting Insight Generation Process using config: {config_file_path}")

    # --- Initialize OpenAI Client ONCE ---
    llm_client = None
    try:
        load_dotenv()
        llm_client = AzureOpenAI(api_key=os.getenv("AZURE_OPENAI_API_KEY"), azure_endpoint=os.getenv("AZURE_OPENAI_API_ENDOINT"), api_version="2024-02-15-preview")
        llm_client.models.list(); print("OpenAI client initialized successfully.")
    except Exception as e: print(f"FATAL ERROR: Could not initialize OpenAI client. Error: {e}"); return

    # 1. Load Unified Config
    mapping_data = {}
    insight_placeholders_map = {}
    summary_placeholders_map = {}
    llm_model_id = "gpt-4o"
    llm_temperature = 0.2
    try:
        print(f"Attempting to load config from: {os.path.abspath(config_file_path)}")
        with open(config_file_path, 'r') as f: mapping_data = json.load(f, strict=False)
        print("--- Successfully loaded JSON data ---")

        # PPTX_FILE_PATH = mapping_data.get("ppt_path")
        # PPTX_OUTPUT_PATH = mapping_data.get("ppt_output_path")
        PPTX_FILE_PATH = mapping_data.get("ppt_output_path")
        PPTX_OUTPUT_PATH = mapping_data.get("ppt_insights_output_path")        
        llm_model_id = mapping_data.get("llm_model_id", llm_model_id)
        llm_temperature = mapping_data.get("llm_temperature", llm_temperature)

        if not PPTX_FILE_PATH: raise ValueError("'ppt_path' not found in config.")
        if not PPTX_OUTPUT_PATH: raise ValueError("'ppt_output_path' not found in config.")

        config_dir = os.path.dirname(config_file_path)
        if not os.path.isabs(PPTX_FILE_PATH): PPTX_FILE_PATH = os.path.join(config_dir, PPTX_FILE_PATH)
        if not os.path.isabs(PPTX_OUTPUT_PATH): PPTX_OUTPUT_PATH = os.path.join(config_dir, PPTX_OUTPUT_PATH)

        print(f"Using PPTX input file: {PPTX_FILE_PATH}")
        print(f"Final PPT output will be saved to: {PPTX_OUTPUT_PATH}")
        print(f"Using LLM Model Deployment: {llm_model_id}")
        print(f"Using LLM Temperature: {llm_temperature}")

        default_excel = mapping_data.get("excel_path")
        if default_excel and not os.path.isabs(default_excel): mapping_data["excel_path"] = os.path.join(config_dir, default_excel)

        insight_placeholders_map = mapping_data.get("insight_placeholders", {})
        summary_placeholders_map = mapping_data.get("slide_summary_placeholders", {})
        print(f"DEBUG: Loaded insight_placeholders_map: {insight_placeholders_map}")
        print(f"DEBUG: Loaded slide_summary_placeholders_map: {summary_placeholders_map}")
        if not insight_placeholders_map: print("Warning: 'insight_placeholders' key not found/empty. Key Highlights will not be pasted.")
        if not summary_placeholders_map: print("Warning: 'slide_summary_placeholders' key not found/empty. Summary Phrases will not be pasted.")

    except FileNotFoundError: print(f"FATAL ERROR: Config file not found at '{os.path.abspath(config_file_path)}'."); return
    except json.JSONDecodeError as e: print(f"FATAL ERROR: Invalid JSON syntax in '{config_file_path}'. Details: {e}"); return
    except Exception as e: print(f"FATAL ERROR: Error loading config JSON '{config_file_path}': {e}"); return

    # 2. Load Presentation
    presentation = None
    try:
        presentation = Presentation(PPTX_FILE_PATH)
        print(f"Opened presentation with {len(presentation.slides)} slides.")
    except Exception as e: print(f"FATAL ERROR: Error opening PPTX file '{PPTX_FILE_PATH}': {e}"); return

    # 3. Initialize Data Structures
    all_extracted_data = {}
    slide_2_context_snippets = {'organic_os': None, 'owned_campaign': None, 'paid_partner': None}
    all_insights_to_paste = {}
    all_prompts_generated = {} # <<< Dictionary to store prompts

    # 4. Define Processing Order
    num_slides_in_ppt = len(presentation.slides)
    all_slide_nums = list(range(1, num_slides_in_ppt + 1))
    slides_to_process = sorted([s for s in all_slide_nums if s in [3, 4, 5]]) \
                      + sorted([s for s in all_slide_nums if s == 2]) \
                      + sorted([s for s in all_slide_nums if s not in [2, 3, 4, 5]])
    print(f"Processing slides in order: {slides_to_process}")

    # 5. Main Processing Loop (Generate Insights)
    for slide_num in slides_to_process:
        if slide_num < 1 or slide_num > num_slides_in_ppt: continue
        print(f"\n===== Processing Slide {slide_num} =====")

        # --- Extract Data Once ---
        if slide_num not in all_extracted_data:
            print(f"  Extracting data for slide {slide_num}...")
            current_slide_data = {}
            try:
                function_1_extract_pptx(slide_num, presentation, current_slide_data)
                function_2_add_chart_data(slide_num, mapping_data, current_slide_data)
                all_extracted_data[slide_num] = current_slide_data.copy()
                print(f"  Data extraction complete for slide {slide_num}.")
            except Exception as e:
                 print(f"  ERROR extracting data for slide {slide_num}: {e}. Skipping insight generation.")
                 continue

        slide_data_for_prompt = all_extracted_data[slide_num]
        # if slide_num == 8:
        #     print(f"\nSLIDE 08 DATA IG: {all_extracted_data[slide_num]}\n")
        # print(f"\nThere it is: {all_extracted_data}\n")
        if slide_num not in all_insights_to_paste: all_insights_to_paste[slide_num] = {}

        # --- Special Handling for Slides 7 and 10 ---
        if slide_num in [7, 10]:
            print(f"  Handling chart-specific insights for slide {slide_num}...")
            chart_textbox_map = {}
            if slide_num == 7:
                chart_textbox_map = {"Chart 9": "Content Placeholder 14", "Chart 2": "Content Placeholder 15"}
            elif slide_num == 10:
                 chart_textbox_map = {"chart_04_02_Instance1": "Content Placeholder 14", "chart_04_02_Instance2": "Content Placeholder 15"} # VERIFY KEYS

            charts_on_slide = slide_data_for_prompt.get('charts_excel_data', [])
            chart_identifiers_from_f1 = slide_data_for_prompt.get('chart_identifiers', [])
            processed_chart_ids = set()

            for chart_id_key, textbox_name in chart_textbox_map.items():
                print(f"    Processing pair: Chart Key='{chart_id_key}' -> Textbox='{textbox_name}'")
                # Debug
                # print(f"\nHi01\nand {charts_on_slide}")
                chart_data_to_send = None; chart_title = chart_id_key; found_match = False
                for i, chart_info in enumerate(charts_on_slide):
                    # print("\nHi02\n")
                    chart_f1_info = chart_identifiers_from_f1[i] if i < len(chart_identifiers_from_f1) else {}
                    chart_name_f1 = chart_f1_info.get('name'); chart_title_f1 = chart_f1_info.get('title')
                    chart_excel_id = chart_info.get('identifier')
                    # match_condition = (chart_excel_id == chart_id_key or chart_title_f1 == chart_id_key) # Adjust matching if needed
                    match_condition = True
                    # print(f"i is: {i}\n and pci is: {processed_chart_ids}\n")
                    if match_condition and i not in processed_chart_ids:
                        # print("\nhi\n")
                        chart_data_to_send = {'name': chart_name_f1 or chart_excel_id, 'title': chart_title_f1 or chart_info.get('mapped_title', chart_id_key), 'data': chart_info.get('data')}
                        chart_title = chart_data_to_send['title'] or chart_data_to_send['name']; print(f"      Found data for chart '{chart_title}'.")
                        processed_chart_ids.add(i); found_match = True; break
                if not found_match or not chart_data_to_send or not chart_data_to_send.get('data'):
                    print(f"      Warning: Data for chart key '{chart_id_key}' not found or empty. Skipping LLM call.")
                    all_insights_to_paste[slide_num][textbox_name] = {'text': "[Data not found for chart]", 'type': 'chart_specific'}; continue

                prompt = build_prompt_for_slide(slide_num, None, chart_focus=chart_data_to_send)
                prompt_key = f"slide_{slide_num}_chart_{chart_data_to_send.get('name','unknown')}" # <<< Store prompt
                all_prompts_generated[prompt_key] = prompt # <<< Store prompt

                if not prompt:
                    print(f"      Skipping LLM call for chart '{chart_title}' as no prompt was generated.")
                    all_insights_to_paste[slide_num][textbox_name] = {'text': "[Prompt generation failed]", 'type': 'chart_specific'}; continue

                print(f"      Calling LLM for chart '{chart_title}'...")
                raw_response_text = function_4_call_llm(prompt, llm_client, model_id=llm_model_id, temperature=llm_temperature)
                parsed_result = parse_llm_response(slide_num, raw_response_text, is_chart_specific=True)
                chart_insight_text = parsed_result['main_insight']
                all_insights_to_paste[slide_num][textbox_name] = {'text': chart_insight_text, 'type': 'chart_specific'}
                print(f"      Stored insight for chart '{chart_title}' targeting textbox '{textbox_name}'.")

            # After processing chart-specific textboxes, create a combined summary for slide 7
            if slide_num == 7:
                sp_shape = summary_placeholders_map.get(str(slide_num))
                if sp_shape:

                    chart_texts = []
                    for k, v in all_insights_to_paste.get(slide_num, {}).items():
                        if isinstance(v, dict) and v.get('type') == 'chart_specific':
                            t = v.get('text') or ''
                            first_line = next((ln.strip() for ln in t.splitlines() if ln.strip()), '')
                            if first_line:
                                # Remove leading numbering/bullets (e.g., '1. ', '- ', '* ')
                                cleaned = re.sub(r'^\s*[\d\-\*]+\.?\s*', '', first_line)
                                chart_texts.append(cleaned)

                    if chart_texts:
                        # Build a prompt from the chart-level insights and ask the LLM
                        joined = '; '.join(chart_texts)
                        prompt = (
                            "You are an expert analyst. Based ONLY on the following chart-level insights, \n"
                            "generate a concise Summary Phrase (maximum 10 words) that captures the most impactful highlight.\n"
                            "Format the output exactly as: Summary Phrase: [Concise phrase]\n\n"
                            f"Chart Insights:\n{joined}"
                        )
                        prompt_key = f"slide_{slide_num}_summary_from_charts"
                        all_prompts_generated[prompt_key] = prompt

                        raw_response = function_4_call_llm(prompt, llm_client, model_id=llm_model_id, temperature=llm_temperature, max_tokens_response=60)
                        parsed = parse_llm_response(slide_num, raw_response)
                        summary_phrase = parsed.get('summary_phrase') or "[No summary generated]"
                    else:
                        summary_phrase = "[No chart insights available to summarize]"

                    all_insights_to_paste[slide_num][sp_shape] = {'text': summary_phrase, 'type': 'summary_phrase'}
                    print(f"    Added generated summary to '{sp_shape}' for slide {slide_num}.")

        # --- Standard Handling for Other Slides ---
        else:
            context_for_slide_2 = slide_2_context_snippets if slide_num == 2 else None
            prompt = build_prompt_for_slide(slide_num, slide_data_for_prompt, context_for_slide_2)
            prompt_key = f"slide_{slide_num}_full" # <<< Store prompt
            all_prompts_generated[prompt_key] = prompt # <<< Store prompt
            # if slide_num == 8:
            #     print("\nWAIT\n")
            #     print(prompt_key, all_prompts_generated[prompt_key])

            if not prompt:
                kh_shape = insight_placeholders_map.get(str(slide_num)); sp_shape = summary_placeholders_map.get(str(slide_num))
                if kh_shape: all_insights_to_paste[slide_num][kh_shape] = {'text': "[No prompt generated]", 'type': 'key_highlight'}
                if sp_shape: all_insights_to_paste[slide_num][sp_shape] = {'text': "[No prompt generated]", 'type': 'summary_phrase'}
                continue

            print(f"  Calling LLM for slide {slide_num}...")
            raw_response_text = function_4_call_llm(prompt, llm_client, model_id=llm_model_id, temperature=llm_temperature)
            parsed_result = parse_llm_response(slide_num, raw_response_text)
            # if slide_num == 8:
            #     print("\nFINAL\n")
            #     print(prompt)
            #     print("\n")
            #     print(parsed_result)
            #     print("\n")
            kh_shape = insight_placeholders_map.get(str(slide_num)); sp_shape = summary_placeholders_map.get(str(slide_num))
            if kh_shape: all_insights_to_paste[slide_num][kh_shape] = {'text': parsed_result['main_insight'], 'type': 'key_highlight'}
            if sp_shape: all_insights_to_paste[slide_num][sp_shape] = {'text': parsed_result['summary_phrase'], 'type': 'summary_phrase'}

            print(f"  Stored parsed result for slide {slide_num}.")
            if slide_num == 3 and parsed_result['context_snippet']: slide_2_context_snippets['organic_os'] = parsed_result['context_snippet']; print(f"    Stored 'organic_os' context for Slide 2: '{parsed_result['context_snippet']}'")
            elif slide_num == 4 and parsed_result['context_snippet']: slide_2_context_snippets['owned_campaign'] = parsed_result['context_snippet']; print(f"    Stored 'owned_campaign' context for Slide 2: '{parsed_result['context_snippet']}'")
            elif slide_num == 5 and parsed_result['context_snippet']: slide_2_context_snippets['paid_partner'] = parsed_result['context_snippet']; print(f"    Stored 'paid_partner' context for Slide 2: '{parsed_result['context_snippet']}'")

    # --- Optional: Print or save the stored prompts ---
    print("\n--- Generated Prompts ---")
    for key, prompt_text in all_prompts_generated.items():
        print(f"\n--- Prompt Key: {key} ---")
        print(prompt_text[:1000] + "...") # Print start of each prompt
    # You could also save `all_prompts_generated` to a file here if needed

    # 6. Paste Insights into PPT using the consolidated map
    if all_insights_to_paste:
        if 'presentation' in locals() and presentation:
             paste_all_text_to_ppt(presentation, all_insights_to_paste)
        else: print("\nError: Presentation object not loaded correctly. Cannot paste insights.")
    else: print("\nNo insights generated or mapped for pasting.")

    # 7. Save Output Files
    print("\n--- Saving Final Outputs ---")
    try:
        if 'presentation' in locals() and presentation:
            presentation.save(PPTX_OUTPUT_PATH)
            print(f"Successfully saved presentation with insights to: {PPTX_OUTPUT_PATH}")
        else: print("Error: Presentation object not loaded correctly. Cannot save.")
    except Exception as e: print(f"ERROR saving modified PowerPoint presentation to '{PPTX_OUTPUT_PATH}': {e}")

    print("\n--- Insight Generation and Pasting Process Complete ---")

# ==============================================================================
# SCRIPT EXECUTION POINT
# ==============================================================================
if __name__ == "__main__":
    # Define the path to your configuration file
    CONFIG_FILE = "ppt_config_mod.json"

    # Call the main logic function
    generate_and_paste_insights(CONFIG_FILE)

